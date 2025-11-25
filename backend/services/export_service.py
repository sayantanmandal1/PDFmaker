"""
Export service for generating Word and PowerPoint documents.
"""
from typing import Optional
from sqlalchemy.orm import Session
from sqlalchemy.exc import SQLAlchemyError
from fastapi import HTTPException, status
from io import BytesIO
import uuid
import logging

from docx import Document
from docx.shared import Pt, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from PIL import Image as PILImage

from models.project import Project
from models.section import Section
from models.slide import Slide
from services.content_service import ContentService
from services.styling_service import StylingService
from services.image_service import ImageService

logger = logging.getLogger(__name__)


class ExportService:
    """Service for exporting projects as Word or PowerPoint documents."""

    @staticmethod
    def export_word_document(
        db: Session,
        project_id: uuid.UUID
    ) -> tuple[BytesIO, str]:
        """
        Export a Word project as a .docx file.
        
        Args:
            db: Database session
            project_id: UUID of the project to export
            
        Returns:
            Tuple of (BytesIO file stream, filename)
            
        Raises:
            HTTPException: If export fails or content is missing
        """
        try:
            # Retrieve project
            project = db.query(Project).filter(Project.id == project_id).first()
            
            if not project:
                raise HTTPException(
                    status_code=status.HTTP_404_NOT_FOUND,
                    detail="Project not found"
                )
            
            if project.document_type != "word":
                raise HTTPException(
                    status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
                    detail="Cannot export non-Word project as Word document"
                )
            
            # Retrieve all sections ordered by position
            sections = ContentService.get_project_sections(db, project_id)
            
            if not sections:
                raise HTTPException(
                    status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
                    detail="Cannot export project without configured sections"
                )
            
            # Check if content has been generated
            if any(section.content is None or section.content.strip() == "" for section in sections):
                raise HTTPException(
                    status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
                    detail="Cannot export project without generated content"
                )
            
            # Create Word document
            doc = ExportService._format_word_document(project, sections)
            
            # Save to BytesIO
            file_stream = BytesIO()
            doc.save(file_stream)
            file_stream.seek(0)
            
            # Generate filename
            filename = f"{project.name.replace(' ', '_')}.docx"
            
            return file_stream, filename
            
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail="Document export failed"
            )

    @staticmethod
    def export_powerpoint_presentation(
        db: Session,
        project_id: uuid.UUID
    ) -> tuple[BytesIO, str]:
        """
        Export a PowerPoint project as a .pptx file.
        
        Args:
            db: Database session
            project_id: UUID of the project to export
            
        Returns:
            Tuple of (BytesIO file stream, filename)
            
        Raises:
            HTTPException: If export fails or content is missing
        """
        try:
            # Retrieve project
            project = db.query(Project).filter(Project.id == project_id).first()
            
            if not project:
                raise HTTPException(
                    status_code=status.HTTP_404_NOT_FOUND,
                    detail="Project not found"
                )
            
            if project.document_type != "powerpoint":
                raise HTTPException(
                    status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
                    detail="Cannot export non-PowerPoint project as PowerPoint presentation"
                )
            
            # Retrieve all slides ordered by position
            slides = ContentService.get_project_slides(db, project_id)
            
            if not slides:
                raise HTTPException(
                    status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
                    detail="Cannot export project without configured slides"
                )
            
            # Check if content has been generated
            if any(slide.content is None or slide.content.strip() == "" for slide in slides):
                raise HTTPException(
                    status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
                    detail="Cannot export project without generated content"
                )
            
            # Create PowerPoint presentation
            prs = ExportService._format_powerpoint_presentation(project, slides)
            
            # Save to BytesIO
            file_stream = BytesIO()
            prs.save(file_stream)
            file_stream.seek(0)
            
            # Generate filename
            filename = f"{project.name.replace(' ', '_')}.pptx"
            
            return file_stream, filename
            
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail="Document export failed"
            )

    @staticmethod
    def _format_word_document(project: Project, sections: list[Section]) -> Document:
        """
        Format sections into a Word document with professional styling and images.
        
        Args:
            project: Project object
            sections: List of Section objects ordered by position
            
        Returns:
            Formatted Document object
        """
        doc = Document()
        
        # Set document properties
        core_properties = doc.core_properties
        core_properties.title = project.name
        core_properties.subject = project.topic
        
        # Apply professional styling to the document
        try:
            StylingService.apply_word_styles(doc)
            logger.info("Applied professional styling to Word document")
        except Exception as e:
            logger.warning(f"Failed to apply Word styles, using defaults: {e}")
        
        # Add document title
        title = doc.add_heading(project.name, level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Add topic as subtitle
        topic_para = doc.add_paragraph(project.topic)
        topic_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        topic_para.runs[0].italic = True
        
        # Add spacing
        doc.add_paragraph()
        
        # Initialize ImageService for downloading images
        image_service = ImageService()
        
        # Add each section
        for section in sections:
            # Add section header as Heading 1
            heading = doc.add_heading(section.header, level=1)
            
            # Apply heading formatting
            try:
                StylingService.format_word_heading(heading, level=1)
            except Exception as e:
                logger.warning(f"Failed to format heading for section '{section.header}': {e}")
            
            # Add section content as body text
            if section.content:
                # Split content into paragraphs
                paragraphs = section.content.split('\n')
                for para_text in paragraphs:
                    if para_text.strip():  # Only add non-empty paragraphs
                        para = doc.add_paragraph(para_text.strip())
                        
                        # Apply body text formatting
                        try:
                            StylingService.format_word_body(para)
                        except Exception as e:
                            logger.warning(f"Failed to format body paragraph: {e}")
            
            # Add image if available
            if section.image_url:
                try:
                    ExportService._add_image_to_word(
                        doc, 
                        section, 
                        image_service
                    )
                    logger.info(f"Added image to section '{section.header}'")
                except Exception as e:
                    logger.warning(f"Failed to add image to section '{section.header}': {e}")
            
            # Add spacing between sections
            doc.add_paragraph()
        
        # Close image service
        image_service.close()
        
        return doc

    @staticmethod
    def _format_powerpoint_presentation(project: Project, slides: list[Slide]) -> Presentation:
        """
        Format slides into a PowerPoint presentation with professional themes and images.
        
        Args:
            project: Project object
            slides: List of Slide objects ordered by position
            
        Returns:
            Formatted Presentation object
        """
        prs = Presentation()
        
        # Set slide dimensions (standard 16:9)
        prs.slide_width = PptxInches(10)
        prs.slide_height = PptxInches(7.5)
        
        # Initialize ImageService for downloading images
        image_service = ImageService()
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
        title_slide = prs.slides.add_slide(title_slide_layout)
        
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        
        title.text = project.name
        subtitle.text = project.topic
        
        # Apply font sizing to title slide
        try:
            StylingService.set_slide_fonts(title_slide, title_size=44, content_size=28)
        except Exception as e:
            logger.warning(f"Failed to set title slide fonts: {e}")
        
        # Add content slides
        for slide_data in slides:
            # Use title and content layout
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Set slide title
            title_shape = slide.shapes.title
            title_shape.text = slide_data.title
            
            # Add content to the content placeholder
            if slide_data.content:
                # Get the content placeholder (usually index 1)
                content_placeholder = slide.placeholders[1]
                text_frame = content_placeholder.text_frame
                text_frame.clear()  # Clear default text
                
                # Split content into paragraphs
                paragraphs = slide_data.content.split('\n')
                
                for i, para_text in enumerate(paragraphs):
                    if para_text.strip():  # Only add non-empty paragraphs
                        if i == 0:
                            # Use the first paragraph in the existing text frame
                            p = text_frame.paragraphs[0]
                        else:
                            # Add new paragraphs
                            p = text_frame.add_paragraph()
                        
                        p.text = para_text.strip()
                        p.level = 0
                        
                        # Set font size
                        for run in p.runs:
                            run.font.size = PptxPt(20)
            
            # Apply font sizing to content slide
            try:
                StylingService.set_slide_fonts(slide, title_size=40, content_size=20)
            except Exception as e:
                logger.warning(f"Failed to set slide fonts for '{slide_data.title}': {e}")
            
            # Add image if available
            if slide_data.image_url:
                try:
                    ExportService._add_image_to_slide(
                        slide,
                        slide_data,
                        image_service
                    )
                    logger.info(f"Added image to slide '{slide_data.title}'")
                except Exception as e:
                    logger.warning(f"Failed to add image to slide '{slide_data.title}': {e}")
        
        # Apply professional theme to the entire presentation
        try:
            StylingService.apply_powerpoint_theme(prs, theme_name="professional_blue")
            logger.info("Applied professional theme to PowerPoint presentation")
        except Exception as e:
            logger.warning(f"Failed to apply PowerPoint theme: {e}. Using default styling.")
        
        # Close image service
        image_service.close()
        
        return prs
    
    @staticmethod
    def _add_image_to_word(
        document: Document,
        section: Section,
        image_service: ImageService
    ) -> None:
        """
        Add image to Word document with appropriate text wrapping.
        
        Args:
            document: Document object to add image to
            section: Section object containing image metadata
            image_service: ImageService instance for downloading images
            
        Raises:
            Exception: If image download or embedding fails
        """
        if not section.image_url:
            logger.warning(f"No image URL for section '{section.header}'")
            return
        
        try:
            # Download the image
            logger.info(f"Downloading image from {section.image_url}")
            image_bytes = image_service.download_image(section.image_url)
            
            if not image_bytes:
                logger.warning(f"Failed to download image from {section.image_url}")
                return
            
            # Optimize image for Word document
            optimized_bytes = image_service.optimize_for_document(image_bytes, 'word')
            
            if not optimized_bytes:
                logger.warning(f"Failed to optimize image for section '{section.header}'")
                # Try to use original image
                optimized_bytes = image_bytes
            
            # Create BytesIO object for the image
            image_stream = BytesIO(optimized_bytes)
            
            # Get image dimensions to calculate appropriate display size
            try:
                pil_image = PILImage.open(BytesIO(optimized_bytes))
                img_width, img_height = pil_image.size
                
                # Calculate display size (max 6 inches width for inline, 3.5 inches for wrapped)
                if section.image_placement == 'wrapped':
                    max_width_inches = 3.5
                else:
                    # Default to inline
                    max_width_inches = 6.0
                
                # Calculate height maintaining aspect ratio
                aspect_ratio = img_height / img_width
                display_width = Inches(max_width_inches)
                display_height = Inches(max_width_inches * aspect_ratio)
                
                # Limit height to reasonable maximum (8 inches)
                max_height = Inches(8)
                if display_height > max_height:
                    display_height = max_height
                    display_width = Inches(8 / aspect_ratio)
                
            except Exception as e:
                logger.warning(f"Failed to get image dimensions: {e}. Using default size.")
                display_width = Inches(4)
                display_height = Inches(3)
            
            # Add the image to the document
            # Create a new paragraph for the image
            image_paragraph = document.add_paragraph()
            run = image_paragraph.add_run()
            
            # Add the picture to the run
            picture = run.add_picture(image_stream, width=display_width, height=display_height)
            
            # Set text wrapping if specified
            if section.image_placement == 'wrapped':
                try:
                    # Set inline shape to allow text wrapping
                    # Note: python-docx has limited support for text wrapping
                    # The image will be inline by default, which is appropriate for most cases
                    image_paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
                    logger.debug(f"Set image alignment for wrapped placement")
                except Exception as e:
                    logger.warning(f"Failed to set text wrapping: {e}")
            else:
                # Center inline images
                image_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            logger.info(f"Successfully added image to section '{section.header}'")
            
        except Exception as e:
            logger.error(f"Error adding image to Word document: {e}")
            raise
    
    @staticmethod
    def _add_image_to_slide(
        slide,
        slide_data: Slide,
        image_service: ImageService
    ) -> None:
        """
        Add image to PowerPoint slide with appropriate positioning based on placement metadata.
        
        Args:
            slide: Slide object to add image to
            slide_data: Slide model object containing image metadata
            image_service: ImageService instance for downloading images
            
        Raises:
            Exception: If image download or embedding fails
        """
        if not slide_data.image_url:
            logger.warning(f"No image URL for slide '{slide_data.title}'")
            return
        
        try:
            # Download the image
            logger.info(f"Downloading image from {slide_data.image_url}")
            image_bytes = image_service.download_image(slide_data.image_url)
            
            if not image_bytes:
                logger.warning(f"Failed to download image from {slide_data.image_url}")
                return
            
            # Optimize image for PowerPoint
            optimized_bytes = image_service.optimize_for_document(image_bytes, 'powerpoint')
            
            if not optimized_bytes:
                logger.warning(f"Failed to optimize image for slide '{slide_data.title}'")
                # Try to use original image
                optimized_bytes = image_bytes
            
            # Create BytesIO object for the image
            image_stream = BytesIO(optimized_bytes)
            
            # Get image dimensions
            try:
                pil_image = PILImage.open(BytesIO(optimized_bytes))
                img_width, img_height = pil_image.size
                aspect_ratio = img_height / img_width
            except Exception as e:
                logger.warning(f"Failed to get image dimensions: {e}. Using default aspect ratio.")
                aspect_ratio = 0.75  # Default 4:3 aspect ratio
            
            # Determine placement strategy
            placement = slide_data.image_placement or 'foreground'
            position = slide_data.image_position or 'center'
            
            # Get slide dimensions
            slide_width = slide.width
            slide_height = slide.height
            
            if placement == 'background':
                # Add as background image (full slide)
                try:
                    # Position image to cover entire slide
                    left = PptxInches(0)
                    top = PptxInches(0)
                    width = slide_width
                    height = slide_height
                    
                    # Add picture at the back
                    picture = slide.shapes.add_picture(
                        image_stream,
                        left,
                        top,
                        width=width,
                        height=height
                    )
                    
                    # Move to back so text appears on top
                    # Note: python-pptx doesn't have direct z-order control
                    # Images added first appear behind shapes added later
                    # We need to reorder by moving the picture to the back
                    slide.shapes._spTree.remove(picture._element)
                    slide.shapes._spTree.insert(2, picture._element)  # Insert after background
                    
                    logger.info(f"Added background image to slide '{slide_data.title}'")
                    
                except Exception as e:
                    logger.warning(f"Failed to add background image: {e}. Trying foreground placement.")
                    placement = 'foreground'
            
            if placement == 'foreground':
                # Add as foreground image (positioned within content area)
                try:
                    # Calculate size and position based on position parameter
                    # Leave space for text content
                    
                    if position == 'right':
                        # Position on right side
                        width = PptxInches(3.5)
                        height = width * aspect_ratio
                        left = slide_width - width - PptxInches(0.5)
                        top = PptxInches(2)
                        
                    elif position == 'left':
                        # Position on left side
                        width = PptxInches(3.5)
                        height = width * aspect_ratio
                        left = PptxInches(0.5)
                        top = PptxInches(2)
                        
                    elif position == 'bottom':
                        # Position at bottom center
                        width = PptxInches(4)
                        height = width * aspect_ratio
                        left = (slide_width - width) / 2
                        top = slide_height - height - PptxInches(0.5)
                        
                    else:
                        # Default: center position
                        width = PptxInches(4.5)
                        height = width * aspect_ratio
                        
                        # Limit height to avoid covering too much content
                        max_height = PptxInches(4)
                        if height > max_height:
                            height = max_height
                            width = height / aspect_ratio
                        
                        left = (slide_width - width) / 2
                        top = (slide_height - height) / 2
                    
                    # Add picture
                    picture = slide.shapes.add_picture(
                        image_stream,
                        left,
                        top,
                        width=width,
                        height=height
                    )
                    
                    logger.info(f"Added foreground image to slide '{slide_data.title}' at position '{position}'")
                    
                except Exception as e:
                    logger.error(f"Failed to add foreground image: {e}")
                    raise
            
        except Exception as e:
            logger.error(f"Error adding image to PowerPoint slide: {e}")
            raise

