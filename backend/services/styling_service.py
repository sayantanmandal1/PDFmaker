"""
Styling service for applying professional themes and styles to documents.
"""
from typing import Dict, Optional
import logging

from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Pt as PptxPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptxRGBColor

logger = logging.getLogger(__name__)


class StylingService:
    """Service for applying professional styling and themes to documents."""
    
    # Theme color definitions
    THEMES = {
        "professional_blue": {
            "primary": (46, 117, 182),      # #2E75B6
            "secondary": (68, 114, 196),    # #4472C4
            "accent": (91, 155, 213),       # #5B9BD5
            "text": (64, 64, 64),           # #404040
            "background": (255, 255, 255),  # #FFFFFF
        },
        "modern_minimal": {
            "primary": (89, 89, 89),        # #595959
            "secondary": (127, 127, 127),   # #7F7F7F
            "accent": (166, 166, 166),      # #A6A6A6
            "text": (64, 64, 64),           # #404040
            "background": (255, 255, 255),  # #FFFFFF
        },
        "creative_vibrant": {
            "primary": (231, 76, 60),       # #E74C3C
            "secondary": (52, 152, 219),    # #3498DB
            "accent": (46, 204, 113),       # #2ECC71
            "text": (44, 62, 80),           # #2C3E50
            "background": (255, 255, 255),  # #FFFFFF
        },
        "classic_formal": {
            "primary": (31, 56, 100),       # #1F3864
            "secondary": (79, 129, 189),    # #4F81BD
            "accent": (192, 0, 0),          # #C00000
            "text": (0, 0, 0),              # #000000
            "background": (255, 255, 255),  # #FFFFFF
        }
    }
    
    # Default theme
    DEFAULT_THEME = "professional_blue"
    
    @staticmethod
    def get_theme_colors(theme_name: str = None) -> Dict[str, tuple]:
        """
        Get color scheme for specified theme.
        
        Args:
            theme_name: Name of the theme (optional, defaults to professional_blue)
            
        Returns:
            Dictionary with color definitions (RGB tuples)
        """
        if theme_name is None or theme_name not in StylingService.THEMES:
            if theme_name is not None:
                logger.warning(f"Invalid theme name '{theme_name}', using default theme")
            theme_name = StylingService.DEFAULT_THEME
        
        return StylingService.THEMES[theme_name]
    
    @staticmethod
    def apply_powerpoint_theme(
        presentation: Presentation,
        theme_name: str = None
    ) -> None:
        """
        Apply professional theme to PowerPoint presentation.
        
        Args:
            presentation: Presentation object to style
            theme_name: Name of the theme to apply (optional)
        """
        try:
            colors = StylingService.get_theme_colors(theme_name)
            
            # Apply theme to all slides
            for slide_idx, slide in enumerate(presentation.slides):
                try:
                    # Set background color (skip title slide to keep it clean)
                    if slide_idx > 0:  # Skip first slide (title slide)
                        background = slide.background
                        fill = background.fill
                        fill.solid()
                        fill.fore_color.rgb = PptxRGBColor(*colors["background"])
                    
                    # Style title shapes
                    if slide.shapes.title and slide.shapes.title.has_text_frame:
                        title_frame = slide.shapes.title.text_frame
                        for paragraph in title_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.color.rgb = PptxRGBColor(*colors["primary"])
                                run.font.bold = True
                                # Ensure font name is set
                                run.font.name = 'Calibri'
                    
                    # Style other text shapes
                    for shape in slide.shapes:
                        if hasattr(shape, "text_frame") and shape.has_text_frame and shape != slide.shapes.title:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.color.rgb = PptxRGBColor(*colors["text"])
                                    # Ensure font name is set
                                    run.font.name = 'Calibri'
                except Exception as slide_error:
                    logger.warning(f"Failed to apply theme to slide {slide_idx}: {slide_error}")
                    continue
            
            logger.info(f"Applied PowerPoint theme: {theme_name or StylingService.DEFAULT_THEME}")
            
        except Exception as e:
            logger.warning(f"Failed to apply PowerPoint theme: {e}. Using default styling.")
    
    @staticmethod
    def set_slide_fonts(
        slide,
        title_size: int = 40,
        content_size: int = 20
    ) -> None:
        """
        Set appropriate font sizes for slide elements.
        
        Args:
            slide: Slide object to style
            title_size: Font size for title (default 40pt)
            content_size: Font size for content (default 20pt)
        """
        try:
            # Set title font size
            if slide.shapes.title:
                title_frame = slide.shapes.title.text_frame
                for paragraph in title_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = PptxPt(title_size)
                        run.font.bold = True
            
            # Set content font sizes
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape != slide.shapes.title:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = PptxPt(content_size)
            
            logger.debug(f"Set slide fonts: title={title_size}pt, content={content_size}pt")
            
        except Exception as e:
            logger.warning(f"Failed to set slide fonts: {e}. Using default sizes.")
    
    @staticmethod
    def apply_word_styles(document: Document, theme_name: str = None) -> None:
        """
        Apply professional styles to Word document.
        
        Args:
            document: Document object to style
            theme_name: Name of the theme to apply (optional)
        """
        try:
            colors = StylingService.get_theme_colors(theme_name)
            
            # Define styles
            styles = document.styles
            
            # Configure Heading 1 style
            if 'Heading 1' in styles:
                heading1 = styles['Heading 1']
                heading1.font.name = 'Calibri'
                heading1.font.size = Pt(16)
                heading1.font.bold = True
                heading1.font.color.rgb = RGBColor(*colors["primary"])
                heading1.paragraph_format.space_before = Pt(12)
                heading1.paragraph_format.space_after = Pt(6)
            
            # Configure Heading 2 style
            if 'Heading 2' in styles:
                heading2 = styles['Heading 2']
                heading2.font.name = 'Calibri'
                heading2.font.size = Pt(14)
                heading2.font.bold = True
                heading2.font.color.rgb = RGBColor(*colors["primary"])
                heading2.paragraph_format.space_before = Pt(10)
                heading2.paragraph_format.space_after = Pt(4)
            
            # Configure Heading 3 style
            if 'Heading 3' in styles:
                heading3 = styles['Heading 3']
                heading3.font.name = 'Calibri'
                heading3.font.size = Pt(12)
                heading3.font.bold = True
                heading3.font.color.rgb = RGBColor(*colors["text"])
                heading3.paragraph_format.space_before = Pt(8)
                heading3.paragraph_format.space_after = Pt(3)
            
            # Configure Normal style
            if 'Normal' in styles:
                normal = styles['Normal']
                normal.font.name = 'Calibri'
                normal.font.size = Pt(11)
                normal.paragraph_format.line_spacing = 1.15
                normal.paragraph_format.space_after = Pt(8)
            
            logger.info(f"Applied Word styles with theme: {theme_name or StylingService.DEFAULT_THEME}")
            
        except Exception as e:
            logger.warning(f"Failed to apply Word styles: {e}. Using default styling.")
    
    @staticmethod
    def format_word_heading(paragraph, level: int = 1, theme_name: str = None) -> None:
        """
        Format a paragraph as a heading with appropriate style.
        
        Args:
            paragraph: Paragraph object to format
            level: Heading level (1, 2, or 3)
            theme_name: Name of the theme to apply (optional)
        """
        try:
            colors = StylingService.get_theme_colors(theme_name)
            
            # Set heading style based on level
            if level == 1:
                font_size = 16
                space_before = 12
                space_after = 6
                color = colors["primary"]
            elif level == 2:
                font_size = 14
                space_before = 10
                space_after = 4
                color = colors["primary"]
            elif level == 3:
                font_size = 12
                space_before = 8
                space_after = 3
                color = colors["text"]
            else:
                logger.warning(f"Invalid heading level {level}, using level 1")
                font_size = 16
                space_before = 12
                space_after = 6
                color = colors["primary"]
            
            # Apply formatting
            for run in paragraph.runs:
                run.font.name = 'Calibri'
                run.font.size = Pt(font_size)
                run.font.bold = True
                run.font.color.rgb = RGBColor(*color)
            
            paragraph.paragraph_format.space_before = Pt(space_before)
            paragraph.paragraph_format.space_after = Pt(space_after)
            
            logger.debug(f"Formatted heading level {level}")
            
        except Exception as e:
            logger.warning(f"Failed to format heading: {e}. Using default formatting.")
    
    @staticmethod
    def format_word_body(paragraph, theme_name: str = None) -> None:
        """
        Format a paragraph as body text with appropriate style.
        
        Args:
            paragraph: Paragraph object to format
            theme_name: Name of the theme to apply (optional)
        """
        try:
            colors = StylingService.get_theme_colors(theme_name)
            
            # Apply body text formatting
            for run in paragraph.runs:
                run.font.name = 'Calibri'
                run.font.size = Pt(11)
                run.font.color.rgb = RGBColor(*colors["text"])
            
            paragraph.paragraph_format.line_spacing = 1.15
            paragraph.paragraph_format.space_after = Pt(8)
            
            logger.debug("Formatted body paragraph")
            
        except Exception as e:
            logger.warning(f"Failed to format body text: {e}. Using default formatting.")
