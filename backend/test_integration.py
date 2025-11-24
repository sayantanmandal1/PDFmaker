"""
Comprehensive integration tests for the AI Document Generator platform.
Tests the complete user flow: register  login  create project  configure  generate  refine  export
"""
import pytest
import requests
import time
import io
from typing import Dict, Optional
from docx import Document
from pptx import Presentation

# Configuration
BASE_URL = "http://localhost:8000"
API_BASE = f"{BASE_URL}/api"

class TestUser:
    """Helper class to manage test user state."""
    def __init__(self, email: str, password: str, name: str):
        self.email = email
        self.password = password
        self.name = name
        self.token: Optional[str] = None
        self.user_id: Optional[str] = None
        self.projects: list = []
    
    def get_headers(self) -> Dict[str, str]:
        """Get authorization headers."""
        if not self.token:
            raise ValueError("User not logged in")
        return {"Authorization": f"Bearer {self.token}"}


class IntegrationTestSuite:
    """Comprehensive integration test suite."""
    
    def __init__(self):
        self.base_url = BASE_URL
        self.api_base = API_BASE
        self.test_users = []
    
    def test_health_check(self):
        """Test that the API is running."""
        print("\n Testing health check...")
        response = requests.get(f"{self.api_base}/health")
        assert response.status_code == 200, f"Health check failed: {response.text}"
        data = response.json()
        assert data["status"] == "healthy", "API is not healthy"
        print(" Health check passed")
    
    def test_user_registration(self, user: TestUser):
        """Test user registration."""
        print(f"\n Testing user registration for {user.email}...")
        response = requests.post(
            f"{self.api_base}/auth/register",
            json={
                "email": user.email,
                "password": user.password,
                "name": user.name
            }
        )
        assert response.status_code in [200, 201], f"Registration failed (status {response.status_code}): {response.text}"
        data = response.json()
        assert "user_id" in data, "No user_id in response"
        user.user_id = data["user_id"]
        print(f" User registered successfully: {user.user_id}")
    
    def test_duplicate_registration(self, user: TestUser):
        """Test that duplicate registration is rejected."""
        print(f"\n Testing duplicate registration for {user.email}...")
        response = requests.post(
            f"{self.api_base}/auth/register",
            json={
                "email": user.email,
                "password": user.password,
                "name": user.name
            }
        )
        assert response.status_code == 409, f"Expected 409 Conflict, got {response.status_code}"
        print(" Duplicate registration correctly rejected")
    
    def test_user_login(self, user: TestUser):
        """Test user login."""
        print(f"\n Testing user login for {user.email}...")
        response = requests.post(
            f"{self.api_base}/auth/login",
            json={
                "email": user.email,
                "password": user.password
            }
        )
        assert response.status_code == 200, f"Login failed: {response.text}"
        data = response.json()
        assert "access_token" in data, "No access_token in response"
        user.token = data["access_token"]
        print(f" User logged in successfully")
    
    def test_invalid_login(self):
        """Test that invalid credentials are rejected."""
        print("\n Testing invalid login...")
        response = requests.post(
            f"{self.api_base}/auth/login",
            json={
                "email": "invalid@example.com",
                "password": "wrongpassword"
            }
        )
        assert response.status_code == 401, f"Expected 401 Unauthorized, got {response.status_code}"
        print(" Invalid login correctly rejected")
    
    def test_get_user_info(self, user: TestUser):
        """Test getting current user info."""
        print(f"\n Testing get user info for {user.email}...")
        response = requests.get(
            f"{self.api_base}/auth/me",
            headers=user.get_headers()
        )
        assert response.status_code == 200, f"Get user info failed: {response.text}"
        data = response.json()
        assert data["email"] == user.email, "Email mismatch"
        assert data["name"] == user.name, "Name mismatch"
        print(" User info retrieved successfully")
    
    def test_unauthorized_access(self):
        """Test that unauthorized access is rejected."""
        print("\n Testing unauthorized access...")
        response = requests.get(f"{self.api_base}/projects")
        assert response.status_code in [401, 403], f"Expected 401 or 403, got {response.status_code}"
        print(" Unauthorized access correctly rejected")
    
    def test_create_word_project(self, user: TestUser) -> str:
        """Test creating a Word document project."""
        print(f"\n Testing Word project creation for {user.email}...")
        response = requests.post(
            f"{self.api_base}/projects",
            headers=user.get_headers(),
            json={
                "name": "Test Word Document",
                "document_type": "word",
                "topic": "The Impact of Artificial Intelligence on Modern Business"
            }
        )
        assert response.status_code == 200, f"Project creation failed: {response.text}"
        data = response.json()
        assert "id" in data, "No project id in response"
        assert data["document_type"] == "word", "Document type mismatch"
        project_id = data["id"]
        user.projects.append(project_id)
        print(f" Word project created: {project_id}")
        return project_id
    
    def test_create_powerpoint_project(self, user: TestUser) -> str:
        """Test creating a PowerPoint project."""
        print(f"\n Testing PowerPoint project creation for {user.email}...")
        response = requests.post(
            f"{self.api_base}/projects",
            headers=user.get_headers(),
            json={
                "name": "Test PowerPoint Presentation",
                "document_type": "powerpoint",
                "topic": "Digital Transformation Strategies for 2025"
            }
        )
        assert response.status_code == 200, f"Project creation failed: {response.text}"
        data = response.json()
        assert "id" in data, "No project id in response"
        assert data["document_type"] == "powerpoint", "Document type mismatch"
        project_id = data["id"]
        user.projects.append(project_id)
        print(f" PowerPoint project created: {project_id}")
        return project_id
    
    def test_list_projects(self, user: TestUser):
        """Test listing user projects."""
        print(f"\n Testing project list for {user.email}...")
        response = requests.get(
            f"{self.api_base}/projects",
            headers=user.get_headers()
        )
        assert response.status_code == 200, f"List projects failed: {response.text}"
        data = response.json()
        assert len(data) >= len(user.projects), "Project count mismatch"
        print(f" Projects listed: {len(data)} projects")
    
    def test_get_project(self, user: TestUser, project_id: str):
        """Test getting a specific project."""
        print(f"\n Testing get project {project_id}...")
        response = requests.get(
            f"{self.api_base}/projects/{project_id}",
            headers=user.get_headers()
        )
        assert response.status_code == 200, f"Get project failed: {response.text}"
        data = response.json()
        assert data["id"] == project_id, "Project ID mismatch"
        print(f" Project retrieved successfully")
        return data
    
    def test_configure_word_project(self, user: TestUser, project_id: str):
        """Test configuring a Word project with sections."""
        print(f"\n Testing Word project configuration...")
        sections = [
            {"header": "Introduction", "position": 0},
            {"header": "Current State of AI", "position": 1},
            {"header": "Business Applications", "position": 2},
            {"header": "Challenges and Opportunities", "position": 3},
            {"header": "Conclusion", "position": 4}
        ]
        response = requests.put(
            f"{self.api_base}/projects/{project_id}/configuration",
            headers=user.get_headers(),
            json={"sections": sections}
        )
        assert response.status_code == 200, f"Configuration failed: {response.text}"
        print(f" Word project configured with {len(sections)} sections")
    
    def test_configure_powerpoint_project(self, user: TestUser, project_id: str):
        """Test configuring a PowerPoint project with slides."""
        print(f"\n Testing PowerPoint project configuration...")
        slides = [
            {"title": "Introduction to Digital Transformation", "position": 0},
            {"title": "Key Technologies", "position": 1},
            {"title": "Implementation Strategies", "position": 2},
            {"title": "Case Studies", "position": 3},
            {"title": "Roadmap for Success", "position": 4}
        ]
        response = requests.put(
            f"{self.api_base}/projects/{project_id}/configuration",
            headers=user.get_headers(),
            json={"slides": slides}
        )
        assert response.status_code == 200, f"Configuration failed: {response.text}"
        print(f" PowerPoint project configured with {len(slides)} slides")
    
    def test_generate_content(self, user: TestUser, project_id: str):
        """Test content generation."""
        print(f"\n Testing content generation for project {project_id}...")
        print(" This may take a while as it calls the LLM API...")
        response = requests.post(
            f"{self.api_base}/projects/{project_id}/generate",
            headers=user.get_headers()
        )
        assert response.status_code == 200, f"Content generation failed: {response.text}"
        data = response.json()
        assert "sections" in data or "slides" in data, "No content in response"
        
        if "sections" in data:
            print(f" Generated content for {len(data['sections'])} sections")
        else:
            print(f" Generated content for {len(data['slides'])} slides")
    
    def test_refine_section(self, user: TestUser, project_id: str):
        """Test refining a section."""
        print(f"\n Testing section refinement...")
        
        # Get project to find sections
        project = self.test_get_project(user, project_id)
        if "sections" not in project or len(project["sections"]) == 0:
            print("  No sections to refine, skipping")
            return
        
        section_id = project["sections"][0]["id"]
        print(f" Refining section {section_id}...")
        
        response = requests.post(
            f"{self.api_base}/sections/{section_id}/refine",
            headers=user.get_headers(),
            json={"prompt": "Make this more concise and professional"}
        )
        assert response.status_code == 200, f"Refinement failed: {response.text}"
        data = response.json()
        assert "content" in data, "No content in refined section"
        print(" Section refined successfully")
    
    def test_refine_slide(self, user: TestUser, project_id: str):
        """Test refining a slide."""
        print(f"\n Testing slide refinement...")
        
        # Get project to find slides
        project = self.test_get_project(user, project_id)
        if "slides" not in project or len(project["slides"]) == 0:
            print("  No slides to refine, skipping")
            return
        
        slide_id = project["slides"][0]["id"]
        print(f" Refining slide {slide_id}...")
        
        response = requests.post(
            f"{self.api_base}/slides/{slide_id}/refine",
            headers=user.get_headers(),
            json={"prompt": "Add more bullet points and make it more engaging"}
        )
        assert response.status_code == 200, f"Refinement failed: {response.text}"
        data = response.json()
        assert "content" in data, "No content in refined slide"
        print(" Slide refined successfully")
    
    def test_add_feedback(self, user: TestUser, project_id: str):
        """Test adding feedback to content."""
        print(f"\n Testing feedback functionality...")
        
        project = self.test_get_project(user, project_id)
        
        if "sections" in project and len(project["sections"]) > 0:
            section_id = project["sections"][0]["id"]
            response = requests.post(
                f"{self.api_base}/sections/{section_id}/feedback",
                headers=user.get_headers(),
                json={"feedback_type": "like"}
            )
            assert response.status_code == 200, f"Feedback failed: {response.text}"
            print(" Feedback added to section")
        
        if "slides" in project and len(project["slides"]) > 0:
            slide_id = project["slides"][0]["id"]
            response = requests.post(
                f"{self.api_base}/slides/{slide_id}/feedback",
                headers=user.get_headers(),
                json={"feedback_type": "dislike"}
            )
            assert response.status_code == 200, f"Feedback failed: {response.text}"
            print(" Feedback added to slide")
    
    def test_add_comment(self, user: TestUser, project_id: str):
        """Test adding comments to content."""
        print(f"\n Testing comment functionality...")
        
        project = self.test_get_project(user, project_id)
        
        if "sections" in project and len(project["sections"]) > 0:
            section_id = project["sections"][0]["id"]
            response = requests.post(
                f"{self.api_base}/sections/{section_id}/comments",
                headers=user.get_headers(),
                json={"comment_text": "This section needs more examples"}
            )
            assert response.status_code == 200, f"Comment failed: {response.text}"
            print(" Comment added to section")
        
        if "slides" in project and len(project["slides"]) > 0:
            slide_id = project["slides"][0]["id"]
            response = requests.post(
                f"{self.api_base}/slides/{slide_id}/comments",
                headers=user.get_headers(),
                json={"comment_text": "Great slide, very informative"}
            )
            assert response.status_code == 200, f"Comment failed: {response.text}"
            print(" Comment added to slide")
    
    def test_export_word_document(self, user: TestUser, project_id: str):
        """Test exporting a Word document."""
        print(f"\n Testing Word document export...")
        response = requests.get(
            f"{self.api_base}/projects/{project_id}/export",
            headers=user.get_headers()
        )
        assert response.status_code == 200, f"Export failed: {response.text}"
        assert response.headers["content-type"] == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        
        # Verify it's a valid Word document
        doc_bytes = io.BytesIO(response.content)
        doc = Document(doc_bytes)
        assert len(doc.paragraphs) > 0, "Document has no content"
        print(f" Word document exported successfully ({len(response.content)} bytes)")
    
    def test_export_powerpoint_presentation(self, user: TestUser, project_id: str):
        """Test exporting a PowerPoint presentation."""
        print(f"\n Testing PowerPoint export...")
        response = requests.get(
            f"{self.api_base}/projects/{project_id}/export",
            headers=user.get_headers()
        )
        assert response.status_code == 200, f"Export failed: {response.text}"
        assert response.headers["content-type"] == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        
        # Verify it's a valid PowerPoint presentation
        ppt_bytes = io.BytesIO(response.content)
        prs = Presentation(ppt_bytes)
        assert len(prs.slides) > 0, "Presentation has no slides"
        print(f" PowerPoint exported successfully ({len(response.content)} bytes, {len(prs.slides)} slides)")
    
    def test_generate_ai_template_word(self, user: TestUser) -> str:
        """Test AI template generation for Word document."""
        print(f"\n Testing AI template generation for Word...")
        
        # Create a new project for template testing
        response = requests.post(
            f"{self.api_base}/projects",
            headers=user.get_headers(),
            json={
                "name": "AI Template Test - Word",
                "document_type": "word",
                "topic": "Climate Change and Renewable Energy"
            }
        )
        assert response.status_code == 200
        project_id = response.json()["id"]
        
        print(" Generating AI template...")
        response = requests.post(
            f"{self.api_base}/projects/{project_id}/generate-template",
            headers=user.get_headers(),
            json={
                "topic": "Climate Change and Renewable Energy",
                "document_type": "word"
            }
        )
        assert response.status_code == 200, f"Template generation failed: {response.text}"
        data = response.json()
        assert "template" in data, "No template in response"
        assert "headers" in data["template"], "No headers in template"
        print(f" AI template generated with {len(data['template']['headers'])} sections")
        return project_id
    
    def test_generate_ai_template_powerpoint(self, user: TestUser) -> str:
        """Test AI template generation for PowerPoint."""
        print(f"\n Testing AI template generation for PowerPoint...")
        
        # Create a new project for template testing
        response = requests.post(
            f"{self.api_base}/projects",
            headers=user.get_headers(),
            json={
                "name": "AI Template Test - PowerPoint",
                "document_type": "powerpoint",
                "topic": "Future of Remote Work"
            }
        )
        assert response.status_code == 200
        project_id = response.json()["id"]
        
        print(" Generating AI template...")
        response = requests.post(
            f"{self.api_base}/projects/{project_id}/generate-template",
            headers=user.get_headers(),
            json={
                "topic": "Future of Remote Work",
                "document_type": "powerpoint"
            }
        )
        assert response.status_code == 200, f"Template generation failed: {response.text}"
        data = response.json()
        assert "template" in data, "No template in response"
        assert "slide_titles" in data["template"], "No slide_titles in template"
        print(f" AI template generated with {len(data['template']['slide_titles'])} slides")
        return project_id
    
    def test_delete_project(self, user: TestUser, project_id: str):
        """Test deleting a project."""
        print(f"\n Testing project deletion...")
        response = requests.delete(
            f"{self.api_base}/projects/{project_id}",
            headers=user.get_headers()
        )
        assert response.status_code == 200, f"Delete failed: {response.text}"
        
        # Verify project is deleted
        response = requests.get(
            f"{self.api_base}/projects/{project_id}",
            headers=user.get_headers()
        )
        assert response.status_code == 404, "Project still exists after deletion"
        print(" Project deleted successfully")
    
    def test_concurrent_users(self):
        """Test multiple users working concurrently."""
        print(f"\n Testing concurrent user operations...")
        
        user1 = TestUser(
            email=f"concurrent1_{int(time.time())}@example.com",
            password="password123",
            name="Concurrent User 1"
        )
        user2 = TestUser(
            email=f"concurrent2_{int(time.time())}@example.com",
            password="password123",
            name="Concurrent User 2"
        )
        
        # Register and login both users
        self.test_user_registration(user1)
        self.test_user_registration(user2)
        self.test_user_login(user1)
        self.test_user_login(user2)
        
        # Create projects for both users
        project1 = self.test_create_word_project(user1)
        project2 = self.test_create_powerpoint_project(user2)
        
        # Verify user isolation - user1 cannot access user2's project
        print(" Testing user isolation...")
        response = requests.get(
            f"{self.api_base}/projects/{project2}",
            headers=user1.get_headers()
        )
        assert response.status_code == 403, "User isolation violated"
        print(" User isolation verified")
        
        print(" Concurrent user operations completed successfully")
    
    def run_complete_word_workflow(self):
        """Run complete Word document workflow."""
        print("\n" + "="*80)
        print(" TESTING COMPLETE WORD DOCUMENT WORKFLOW")
        print("="*80)
        
        user = TestUser(
            email=f"word_test_{int(time.time())}@example.com",
            password="password123",
            name="Word Test User"
        )
        
        self.test_user_registration(user)
        self.test_user_login(user)
        self.test_get_user_info(user)
        
        project_id = self.test_create_word_project(user)
        self.test_configure_word_project(user, project_id)
        self.test_generate_content(user, project_id)
        self.test_refine_section(user, project_id)
        self.test_add_feedback(user, project_id)
        self.test_add_comment(user, project_id)
        self.test_export_word_document(user, project_id)
        
        print("\n WORD DOCUMENT WORKFLOW COMPLETED SUCCESSFULLY")
    
    def run_complete_powerpoint_workflow(self):
        """Run complete PowerPoint workflow."""
        print("\n" + "="*80)
        print(" TESTING COMPLETE POWERPOINT WORKFLOW")
        print("="*80)
        
        user = TestUser(
            email=f"ppt_test_{int(time.time())}@example.com",
            password="password123",
            name="PowerPoint Test User"
        )
        
        self.test_user_registration(user)
        self.test_user_login(user)
        self.test_get_user_info(user)
        
        project_id = self.test_create_powerpoint_project(user)
        self.test_configure_powerpoint_project(user, project_id)
        self.test_generate_content(user, project_id)
        self.test_refine_slide(user, project_id)
        self.test_add_feedback(user, project_id)
        self.test_add_comment(user, project_id)
        self.test_export_powerpoint_presentation(user, project_id)
        
        print("\n POWERPOINT WORKFLOW COMPLETED SUCCESSFULLY")
    
    def run_ai_template_workflow(self):
        """Run AI template generation workflow."""
        print("\n" + "="*80)
        print(" TESTING AI TEMPLATE GENERATION WORKFLOW")
        print("="*80)
        
        user = TestUser(
            email=f"template_test_{int(time.time())}@example.com",
            password="password123",
            name="Template Test User"
        )
        
        self.test_user_registration(user)
        self.test_user_login(user)
        
        word_project = self.test_generate_ai_template_word(user)
        ppt_project = self.test_generate_ai_template_powerpoint(user)
        
        print("\n AI TEMPLATE WORKFLOW COMPLETED SUCCESSFULLY")
    
    def run_all_tests(self):
        """Run all integration tests."""
        print("\n" + "="*80)
        print("STARTING COMPREHENSIVE INTEGRATION TESTS")
        print("="*80)
        
        try:
            # Basic health check
            self.test_health_check()
            
            # Test authentication errors
            self.test_invalid_login()
            self.test_unauthorized_access()
            
            # Run complete workflows
            self.run_complete_word_workflow()
            self.run_complete_powerpoint_workflow()
            self.run_ai_template_workflow()
            
            # Test concurrent users
            self.test_concurrent_users()
            
            print("\n" + "="*80)
            print(" ALL INTEGRATION TESTS PASSED!")
            print("="*80)
            
        except AssertionError as e:
            print(f"\n TEST FAILED: {e}")
            raise
        except Exception as e:
            print(f"\n UNEXPECTED ERROR: {e}")
            raise


if __name__ == "__main__":
    suite = IntegrationTestSuite()
    suite.run_all_tests()

